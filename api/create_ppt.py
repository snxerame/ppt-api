import io
import os
import json
import requests
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

LOGO_PATH = "sp_global_logo.png"  # Ensure this is in your project root, or leave as is to skip logo

def set_a4_landscape(prs):
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(21)

def add_footer_with_logo(prs, slide, page_num):
    logo_width = Cm(6.0)
    logo_height = Cm(2.4)
    footer_y = Cm(18.03)
    try:
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Cm(1), footer_y, width=logo_width, height=logo_height)
            left_text_x = Cm(7.2)
        else:
            left_text_x = Cm(1)
    except Exception:
        left_text_x = Cm(1)
    left_box = slide.shapes.add_textbox(left_text_x, footer_y, Cm(10), Cm(1.5))
    left_frame = left_box.text_frame
    left_frame.clear()
    p_left = left_frame.add_paragraph()
    p_left.text = (
        "Permission to reprint or distribute any content from this presentation requires "
        "the prior written approval of S&P Global Market Intelligence."
    )
    p_left.font.size = Pt(10)
    p_left.font.color.rgb = RGBColor(128, 128, 128)
    p_left.alignment = PP_ALIGN.LEFT
    right_box = slide.shapes.add_textbox(prs.slide_width - Cm(3), footer_y, Cm(2.5), Cm(1.5))
    right_frame = right_box.text_frame
    right_frame.clear()
    p_right = right_frame.add_paragraph()
    p_right.text = str(page_num)
    p_right.font.size = Pt(14)
    p_right.font.color.rgb = RGBColor(128, 128, 128)
    p_right.alignment = PP_ALIGN.RIGHT

def add_dates_available_box(slide, left, top, width, height, dates_text):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_box.fill.solid()
    text_box.fill.fore_color.rgb = RGBColor(224, 234, 238)
    text_box.line.color.rgb = RGBColor(224, 234, 238)
    tf = text_box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run1 = p.add_run()
    run1.text = "Proposed Dates: "
    run1.font.size = Pt(14)
    run1.font.bold = True
    run1.font.color.rgb = RGBColor(204, 0, 0)
    run2 = p.add_run()
    run2.text = dates_text
    run2.font.size = Pt(14)
    run2.font.bold = False
    run2.font.color.rgb = RGBColor(0, 0, 0)
    return text_box

def create_front_page(prs, heading, date_to_present):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(153, 153, 153)
    shape.line.fill.background()
    brand_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(8), Cm(2))
    brand_frame = brand_box.text_frame
    brand_frame.clear()
    p1 = brand_frame.add_paragraph()
    p1.text = "S&P Global"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p2 = brand_frame.add_paragraph()
    p2.text = "Market Intelligence"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    title_box = slide.shapes.add_textbox(Cm(1), Cm(6), Cm(22), Cm(4))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.add_paragraph()
    p.text = heading
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    date_box = slide.shapes.add_textbox(Cm(1), Cm(17), Cm(8), Cm(2))
    date_frame = date_box.text_frame
    date_frame.clear()
    p_date = date_frame.add_paragraph()
    p_date.text = date_to_present
    p_date.font.size = Pt(32)
    p_date.font.color.rgb = RGBColor(255, 255, 255)
    bottom_right_box = slide.shapes.add_textbox(prs.slide_width - Cm(9), prs.slide_height - Cm(2), Cm(8), Cm(1))
    bottom_right_frame = bottom_right_box.text_frame
    bottom_right_frame.clear()
    p_br = bottom_right_frame.add_paragraph()
    p_br.text = "S&P Global Market Intelligence"
    p_br.font.size = Pt(14)
    p_br.font.color.rgb = RGBColor(255, 255, 255)
    p_br.alignment = PP_ALIGN.RIGHT

def create_content_slide(prs, idx, venue_info, page_num):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    venue_name = venue_info.get("venue_name", "")
    venue_city = venue_info.get("venue_city", "")
    venue_guest_rooms = venue_info.get("venue_guest_rooms", "")
    proposed_dates = venue_info.get("proposed_dates", "")
    average_daily_rate = venue_info.get("average_daily_rate", "")
    total_fandb = venue_info.get("total_FandB", "")
    additional_fees = venue_info.get("additional_fees", "")
    heading_val = f"Recommendation #{idx+1} – {venue_name} : {venue_guest_rooms} rooms"
    heading_box = slide.shapes.add_textbox(Cm(1), Cm(1), Cm(21), Cm(2))
    heading_frame = heading_box.text_frame
    heading_frame.clear()
    p_heading = heading_frame.add_paragraph()
    p_heading.text = heading_val
    p_heading.font.size = Pt(32)
    p_heading.font.bold = True
    add_dates_available_box(slide, Cm(1), Cm(5.54), Cm(10), Cm(1.2), proposed_dates)
    overview_text = (
        f"• City: {venue_city}\n"
        f"• Guest Rooms: {venue_guest_rooms}\n"
        f"• Average Daily Rate: {average_daily_rate}\n"
        f"• Total Food & Beverages: {total_fandb}\n"
        f"• Additional Fees: {additional_fees}"
    )
    overview_left = Cm(1)
    overview_top = Cm(8)
    overview_width = Cm(12)
    overview_height = Cm(7)
    box_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        overview_left,
        overview_top,
        overview_width,
        overview_height
    )
    box_shape.line.color.rgb = RGBColor(0, 0, 0)
    box_shape.line.width = Pt(2)
    box_shape.fill.background()
    overview_box = slide.shapes.add_textbox(
        overview_left,
        overview_top,
        overview_width,
        overview_height
    )
    overview_frame = overview_box.text_frame
    overview_frame.clear()
    p_overview = overview_frame.add_paragraph()
    p_overview.text = "Hotel Overview"
    p_overview.font.size = Pt(16)
    p_overview.font.bold = True
    p_overview.font.color.rgb = RGBColor(255, 0, 0)
    p_overview.space_after = Pt(8)
    p_overview2 = overview_frame.add_paragraph()
    p_overview2.text = overview_text
    p_overview2.font.size = Pt(14)
    p_overview2.font.color.rgb = RGBColor(0, 0, 0)
    img_w = Cm(7)
    img_h = Cm(4)
    gap_horizontal = Cm(0.5)
    gap_vertical = Cm(2.54)
    img_left = Cm(14.5)
    img_top = Cm(5.54)
    labels = ["Main Ballroom", "Bedroom", "Breakout room", "Outdoor space"]
    positions = [
        (img_left, img_top),
        (img_left + img_w + gap_horizontal, img_top),
        (img_left, img_top + img_h + gap_vertical),
        (img_left + img_w + gap_horizontal, img_top + img_h + gap_vertical)
    ]
    for label, (left, top) in zip(labels, positions):
        label_top = top + img_h + Cm(0.2)
        label_height = Cm(1.0)
        shape = slide.shapes.add_shape(1, left, top, img_w, img_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        label_box = slide.shapes.add_textbox(left, label_top, img_w, label_height)
        label_frame = label_box.text_frame
        label_frame.clear()
        p_label = label_frame.add_paragraph()
        p_label.text = label
        p_label.font.size = Pt(12)
        p_label.alignment = PP_ALIGN.CENTER
    add_footer_with_logo(prs, slide, page_num)

def parse_input(text):
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    heading = lines[0]
    date_to_present = lines[1]
    num = int(lines[2])
    recs = [{} for _ in range(num)]
    for arg in lines[3:]:
        if '=' in arg:
            key, value = arg.split('=', 1)
            if "." in key:
                prefix, param = key.split(".", 1)
                if prefix.startswith("R") and prefix[1:].isdigit():
                    idx = int(prefix[1:]) - 1
                    if 0 <= idx < num:
                        recs[idx][param] = value
    return heading, date_to_present, num, recs

def ppt_api_main(event_body):
    heading, date_to_present, num, recs = parse_input(event_body)
    prs = Presentation()
    set_a4_landscape(prs)
    create_front_page(prs, heading, date_to_present)
    slide_num = 2
    for idx, rec in enumerate(recs):
        create_content_slide(prs, idx, rec, slide_num)
        slide_num += 1
    ppt_mem = io.BytesIO()
    prs.save(ppt_mem)
    ppt_mem.seek(0)
    VERCEL_BLOB_WRITE_URL = "https://api.vercel.com/v8/blob/upload"
    BLOB_TOKEN = os.environ.get("BLOB_READ_WRITE_TOKEN")
    filename = f"{heading.replace(' ', '_')}.pptx" if heading else "presentation.pptx"
    headers = {"Authorization": f"Bearer {BLOB_TOKEN}"}
    files = {"file": (filename, ppt_mem.read())}
    params = {"access": "public"}
    response = requests.post(VERCEL_BLOB_WRITE_URL, headers=headers, files=files, data=params)
    if response.ok:
        url = response.json().get('url', None)
        return {"url": url}, 200
    else:
        try:
            resp_json = response.json()
        except Exception:
            resp_json = {}
        return {"error": "Upload failed", "detail": resp_json}, 500

def handler(request):
    try:
        if request.method == 'POST':
            try:
                req_json = request.get_json(force=True, silent=True)
            except Exception:
                req_json = None
            if req_json and 'text' in req_json:
                body = req_json["text"]
            else:
                body = request.get_data(as_text=True)
            data, status = ppt_api_main(body)
            return (json.dumps(data), status, {'Content-Type': 'application/json'})
        else:
            return (json.dumps({"error": "Method not allowed"}), 405, {'Content-Type': 'application/json'})
    except Exception as e:
        return (json.dumps({"error": str(e)}), 500, {'Content-Type': 'application/json'})
